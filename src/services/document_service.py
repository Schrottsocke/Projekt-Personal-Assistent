"""
Document Service: Erstellt PowerPoint-Präsentationen (.pptx) und Excel-Tabellen (.xlsx).
Dateien werden in data/documents/ gespeichert und nach dem Senden gelöscht.
"""

import logging
import uuid
from datetime import datetime
from pathlib import Path

from config.settings import settings

logger = logging.getLogger(__name__)

# Maximale Zeilen/Spalten für Markdown-Darstellung im Chat
MARKDOWN_MAX_COLS = 5
MARKDOWN_MAX_ROWS = 15


class DocumentService:
    def __init__(self):
        self._docs_dir: Path = settings.DOCUMENTS_DIR

    async def initialize(self):
        self._docs_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"Document Service initialisiert. Verzeichnis: {self._docs_dir}")

    # ── Tabellen ─────────────────────────────────────────────────────────────

    def create_markdown_table(self, headers: list[str], rows: list[list[str]]) -> str:
        """
        Erzeugt eine Monospace-Tabelle als Telegram-Codeblock.
        Geeignet für kleine Tabellen (≤ MARKDOWN_MAX_COLS Spalten, ≤ MARKDOWN_MAX_ROWS Zeilen).
        """
        if not headers:
            return "_(Leere Tabelle)_"

        # Spaltenbreiten berechnen
        col_widths = [len(str(h)) for h in headers]
        for row in rows:
            for i, cell in enumerate(row):
                if i < len(col_widths):
                    col_widths[i] = max(col_widths[i], len(str(cell)))

        def fmt_row(cells: list) -> str:
            parts = []
            for i, cell in enumerate(cells):
                w = col_widths[i] if i < len(col_widths) else 10
                parts.append(str(cell).ljust(w))
            return "| " + " | ".join(parts) + " |"

        separator = "|" + "|".join("-" * (w + 2) for w in col_widths) + "|"

        lines = [fmt_row(headers), separator]
        for row in rows:
            lines.append(fmt_row(row))

        return "```\n" + "\n".join(lines) + "\n```"

    def create_excel(
        self,
        title: str,
        headers: list[str],
        rows: list[list],
        sheet_name: str = "Tabelle",
    ) -> Path:
        """
        Erstellt eine Excel-Datei (.xlsx) und gibt den Pfad zurück.
        Der Aufrufer ist verantwortlich, die Datei nach dem Senden zu löschen.
        """
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name[:31]  # Excel max 31 Zeichen

        # Titelzeile
        ws.merge_cells(f"A1:{chr(64 + len(headers))}1")
        title_cell = ws["A1"]
        title_cell.value = title
        title_cell.font = Font(bold=True, size=14, color="FFFFFF")
        title_cell.fill = PatternFill("solid", fgColor="1F497D")
        title_cell.alignment = Alignment(horizontal="center", vertical="center")
        ws.row_dimensions[1].height = 24

        # Header-Zeile
        header_fill = PatternFill("solid", fgColor="4472C4")
        header_font = Font(bold=True, color="FFFFFF")
        thin = Side(style="thin", color="AAAAAA")
        border = Border(left=thin, right=thin, top=thin, bottom=thin)

        for col_idx, header in enumerate(headers, start=1):
            cell = ws.cell(row=2, column=col_idx, value=str(header))
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")
            cell.border = border

        # Daten-Zeilen
        alt_fill = PatternFill("solid", fgColor="DCE6F1")
        for row_idx, row_data in enumerate(rows, start=3):
            fill = alt_fill if row_idx % 2 == 1 else PatternFill()
            for col_idx, value in enumerate(row_data, start=1):
                cell = ws.cell(
                    row=row_idx,
                    column=col_idx,
                    value=str(value) if value is not None else "",
                )
                cell.fill = fill
                cell.border = border
                cell.alignment = Alignment(wrap_text=True)

        # Spaltenbreiten automatisch anpassen
        for col_idx, header in enumerate(headers, start=1):
            col_letter = chr(64 + col_idx)
            max_len = len(str(header))
            for row_data in rows:
                if col_idx - 1 < len(row_data):
                    max_len = max(max_len, len(str(row_data[col_idx - 1] or "")))
            ws.column_dimensions[col_letter].width = min(max_len + 4, 40)

        filename = f"tabelle_{uuid.uuid4().hex[:8]}.xlsx"
        path = self._docs_dir / filename
        wb.save(path)
        logger.info(f"Excel erstellt: {path}")
        return path

    # ── Präsentationen ───────────────────────────────────────────────────────

    def create_presentation(
        self,
        title: str,
        slides: list[dict],
    ) -> Path:
        """
        Erstellt eine PowerPoint-Präsentation (.pptx).
        slides: [{"title": str, "bullets": [str, ...]}, ...]
        Der Aufrufer ist verantwortlich, die Datei nach dem Senden zu löschen.
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Farben
        COLOR_DARK = RGBColor(0x1F, 0x49, 0x7D)  # Dunkelblau
        _COLOR_ACCENT = RGBColor(0x44, 0x72, 0xC4)  # Mittelblau
        COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        COLOR_TEXT = RGBColor(0x26, 0x26, 0x26)  # Fast Schwarz
        COLOR_LIGHT = RGBColor(0xDC, 0xE6, 0xF1)  # Hellblau

        slide_layouts = prs.slide_layouts

        # ── Titelfolie ──────────────────────────────────────────────────────
        title_slide = prs.slides.add_slide(slide_layouts[6])  # Blank

        # Hintergrund dunkelblau
        background = title_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_DARK

        # Titeltext
        txBox = title_slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(1.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE

        # Datum
        date_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.33), Inches(0.6))
        dtf = date_box.text_frame
        dp = dtf.paragraphs[0]
        dp.alignment = PP_ALIGN.CENTER
        drun = dp.add_run()
        drun.text = datetime.now().strftime("%d. %B %Y")
        drun.font.size = Pt(18)
        drun.font.color.rgb = COLOR_LIGHT

        # ── Inhaltsfolien ───────────────────────────────────────────────────
        for slide_data in slides:
            slide_title = slide_data.get("title", "")
            bullets = slide_data.get("bullets", [])

            content_slide = prs.slides.add_slide(slide_layouts[6])  # Blank

            # Weißer Hintergrund
            bg = content_slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLOR_WHITE

            # Blauer Balken oben
            bar = content_slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                0,
                0,
                prs.slide_width,
                Inches(1.3),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLOR_DARK
            bar.line.fill.background()

            # Folientitel im blauen Balken
            title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.33), Inches(1.0))
            ttf = title_box.text_frame
            tp = ttf.paragraphs[0]
            trun = tp.add_run()
            trun.text = slide_title
            trun.font.size = Pt(28)
            trun.font.bold = True
            trun.font.color.rgb = COLOR_WHITE

            # Bullet-Points
            if bullets:
                content_box = content_slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.73), Inches(5.5))
                ctf = content_box.text_frame
                ctf.word_wrap = True

                for i, bullet in enumerate(bullets[:8]):  # max 8 Bullets
                    if i == 0:
                        cp = ctf.paragraphs[0]
                    else:
                        cp = ctf.add_paragraph()
                    cp.space_before = Pt(6)
                    crun = cp.add_run()
                    crun.text = f"• {bullet}"
                    crun.font.size = Pt(18)
                    crun.font.color.rgb = COLOR_TEXT

        filename = f"praesentation_{uuid.uuid4().hex[:8]}.pptx"
        path = self._docs_dir / filename
        prs.save(path)
        logger.info(f"Präsentation erstellt: {path} ({len(slides)} Folien)")
        return path

    def is_small_table(self, headers: list, rows: list) -> bool:
        """Gibt True zurück wenn die Tabelle als Codeblock in den Chat passt."""
        return len(headers) <= MARKDOWN_MAX_COLS and len(rows) <= MARKDOWN_MAX_ROWS
